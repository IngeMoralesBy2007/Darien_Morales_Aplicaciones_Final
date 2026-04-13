import os

import io

import json

import re

import asyncio

import tempfile

from typing import Dict, Any, List

from fastapi import FastAPI, UploadFile, File, HTTPException, Form

from fastapi.middleware.cors import CORSMiddleware

from fastapi.responses import Response

from pydantic import BaseModel

import pandas as pd

import httpx

import base64

from dotenv import load_dotenv

from docx import Document

from docx.shared import Inches

from pptx import Presentation

from pptx.util import Pt

from fpdf import FPDF

# ═══ [INYECCIÓN] Matplotlib para gráficos reales en Word ═══

import matplotlib

matplotlib.use('Agg')  # Backend no-interactivo para servidor

import matplotlib.pyplot as plt

# ═══ [INYECCIÓN] SSE para streaming en tiempo real ═══

from sse_starlette.sse import EventSourceResponse

from langchain_google_genai import ChatGoogleGenerativeAI

from langchain_openai import ChatOpenAI

from langchain_groq import ChatGroq

from langchain_experimental.agents.agent_toolkits import create_pandas_dataframe_agent

load_dotenv()

GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")

OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")

GROQ_API_KEY = os.getenv("GROQ_API_KEY")

COHERE_API_KEY = os.getenv("COHERE_API_KEY")

OPENROUTER_API_KEY = os.getenv("OPENROUTER_API_KEY")

app = FastAPI(title="Nexus Analytic API")

# Allow Frontend CORS

app.add_middleware(

    CORSMiddleware,

    allow_origins=["*"],

    allow_credentials=True,

    allow_methods=["*"],

    allow_headers=["*"],

)

# In-Memory State to emulate Streamlit session

class SessionState:

    def __init__(self):

        self.df: pd.DataFrame = None

        self.messages: List[Dict[str, str]] = []

        self.word_bytes: bytes = None

        self.pdf_bytes: bytes = None

        self.pptx_bytes: bytes = None

        # ═══ [INYECCIÓN] Almacenar Excel bytes ═══

        self.excel_bytes: bytes = None

sessions: Dict[str, SessionState] = {}

class ChatRequest(BaseModel):

    session_id: str

    prompt: str

    model_choice: str

def get_session(session_id: str) -> SessionState:

    if session_id not in sessions:

        sessions[session_id] = SessionState()

    return sessions[session_id]

def get_llm(model_choice: str):

    load_dotenv(override=True)

    google_key = os.getenv("GOOGLE_API_KEY")

    openai_key = os.getenv("OPENAI_API_KEY")

    groq_key = os.getenv("GROQ_API_KEY")

    cohere_key = os.getenv("COHERE_API_KEY")

    openrouter_key = os.getenv("OPENROUTER_API_KEY")

    if model_choice == "Google Gemini 1.5 Flash" and google_key:

        return ChatGoogleGenerativeAI(model="gemini-1.5-flash", temperature=0.1, google_api_key=google_key)

    elif model_choice == "ChatGPT (GPT-4o)" and openai_key:

        return ChatOpenAI(model="gpt-4o", temperature=0.1, api_key=openai_key)

    elif model_choice == "Copilot (Azure/OpenAI)" and openai_key:

        return ChatOpenAI(model="gpt-4o", temperature=0.1, api_key=openai_key)

    # ═══ [INYECCIÓN] Groq ya integrado, confirmado ═══

    elif model_choice == "Llama-3 70B (Groq)" and groq_key:

        return ChatGroq(model_name="llama-3.3-70b-versatile", temperature=0.1, groq_api_key=groq_key)

    elif model_choice == "Cohere (Command R)" and cohere_key:

        from langchain_cohere import ChatCohere

        return ChatCohere(model="command-r-plus", temperature=0.1, cohere_api_key=cohere_key)

    elif model_choice == "OpenRouter (Llama Free)" and openrouter_key:

        return ChatOpenAI(

            model="meta-llama/llama-3-8b-instruct:free",

            temperature=0.1,

            api_key=openrouter_key,

            base_url="https://openrouter.ai/api/v1"

        )

    return None

# ═══════════════════════════════════════════════════════════════

# [INYECCIÓN] Generador de gráficos reales con Matplotlib

# ═══════════════════════════════════════════════════════════════

def generar_graficos_matplotlib(df: pd.DataFrame) -> List[io.BytesIO]:

    """Genera gráficos reales del DataFrame y los retorna como imágenes en memoria."""

    graficos = []

    numeric_cols = df.select_dtypes(include='number').columns.tolist()

    # Gráfico 1: Histograma de la primera columna numérica

    if len(numeric_cols) >= 1:

        fig, ax = plt.subplots(figsize=(8, 4))

        ax.hist(df[numeric_cols[0]].dropna(), bins=20, color='#00BCD4', edgecolor='#0B0F19', alpha=0.85)

        ax.set_title(f'Distribución de {numeric_cols[0]}', fontsize=14, fontweight='bold', color='#333')

        ax.set_xlabel(numeric_cols[0])

        ax.set_ylabel('Frecuencia')

        ax.grid(axis='y', alpha=0.3)

        buf = io.BytesIO()

        fig.savefig(buf, format='png', dpi=150, bbox_inches='tight', facecolor='white')

        plt.close(fig)

        buf.seek(0)

        graficos.append(buf)

    # Gráfico 2: Correlación entre las dos primeras columnas numéricas

    if len(numeric_cols) >= 2:

        fig, ax = plt.subplots(figsize=(8, 4))

        ax.scatter(df[numeric_cols[0]], df[numeric_cols[1]], color='#7C4DFF', alpha=0.6, s=20)

        ax.set_title(f'{numeric_cols[0]} vs {numeric_cols[1]}', fontsize=14, fontweight='bold', color='#333')

        ax.set_xlabel(numeric_cols[0])

        ax.set_ylabel(numeric_cols[1])

        ax.grid(alpha=0.3)

        buf = io.BytesIO()

        fig.savefig(buf, format='png', dpi=150, bbox_inches='tight', facecolor='white')

        plt.close(fig)

        buf.seek(0)

        graficos.append(buf)

    # Gráfico 3: Top 10 valores de la primera columna categórica (si existe)

    cat_cols = df.select_dtypes(include=['object', 'category']).columns.tolist()

    if cat_cols:

        top_vals = df[cat_cols[0]].value_counts().head(10)

        fig, ax = plt.subplots(figsize=(8, 4))

        top_vals.plot(kind='barh', ax=ax, color='#00E5FF', edgecolor='#0B0F19')

        ax.set_title(f'Top 10 - {cat_cols[0]}', fontsize=14, fontweight='bold', color='#333')

        ax.set_xlabel('Frecuencia')

        ax.grid(axis='x', alpha=0.3)

        buf = io.BytesIO()

        fig.savefig(buf, format='png', dpi=150, bbox_inches='tight', facecolor='white')

        plt.close(fig)

        buf.seek(0)

        graficos.append(buf)

    return graficos

# ═══════════════════════════════════════════════════════════════

# [INYECCIÓN] Word con gráficos reales insertados

# ═══════════════════════════════════════════════════════════════

def generar_informe_word(texto_analisis, df: pd.DataFrame = None):

    doc = Document()

    doc.add_heading('Informe Analítico Ejecutivo - Nexus', 0)

    doc.add_heading('Resumen de Hallazgos', level=1)

    doc.add_paragraph(texto_analisis)

    # ═══ [INYECCIÓN] Insertar gráficos reales generados por Matplotlib ═══

    if df is not None:

        graficos = generar_graficos_matplotlib(df)

        if graficos:

            doc.add_heading('Visualizaciones de Datos', level=1)

            for i, grafico_buf in enumerate(graficos):

                doc.add_picture(grafico_buf, width=Inches(5.5))

                doc.add_paragraph(f'Figura {i + 1}', style='Caption')

    doc.add_paragraph('\n---\nGenerado automáticamente por Nexus Analytic.', style='Intense Quote')

    docx_io = io.BytesIO()

    doc.save(docx_io)

    return docx_io.getvalue()

def generar_informe_pdf(texto_analisis, df: pd.DataFrame = None):

    pdf = FPDF()

    pdf.add_page()

    pdf.set_font("Arial", 'B', 16)

    pdf.cell(0, 10, "Informe Analitico Ejecutivo", ln=True, align="C")

    pdf.ln(5)

    pdf.set_font("Arial", size=12)

    texto_seguro = texto_analisis.replace('·', '*').replace('©', '(c)').replace('—', '-').replace('–', '-').encode('latin-1', 'replace').decode('latin-1')

    pdf.multi_cell(0, 8, texto_seguro)

    if df is not None:

        graficos = generar_graficos_matplotlib(df)

        if graficos:

            pdf.add_page()

            pdf.set_font("Arial", 'B', 14)

            pdf.cell(0, 10, "Visualizaciones de Datos", ln=True, align="L")

            pdf.ln(5)

            # Agregar imágenes temporales a FPDF

            for grafico_buf in graficos:

                with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tf:

                    tf.write(grafico_buf.getvalue())

                    tmp_img = tf.name

                pdf.image(tmp_img, x=20, w=170)

                import os

                try: os.unlink(tmp_img)

                except: pass

    return bytes(pdf.output(dest='S').encode('latin-1'))

def generar_pptx_presentacion(texto_analisis, llm, df: pd.DataFrame = None):

    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[0])

    slide.shapes.title.text = "Informe Ejecutivo"

    prompt = f"Estructura 3 diapositivas. Responde con un array JSON: 'titulo', 'puntos', 'notas_orador'. Reporte: {texto_analisis}"

    fallback = [{"titulo": "Conclusiones Principales", "puntos": ["Revisar Word."], "notas_orador": ""}]

    datos_diapositivas = fallback

    if llm:

        try:

            respuesta_cruda = llm.invoke(prompt).content

            match = re.search(r'\[\s*\{.*\}\s*\]', respuesta_cruda, re.DOTALL)

            if match:

                datos_diapositivas = json.loads(match.group(0))

            else:

                datos_diapositivas = json.loads(respuesta_cruda)

        except: pass

    for info in datos_diapositivas:

        try:

            sld = prs.slides.add_slide(prs.slide_layouts[1])

            sld.shapes.title.text = str(info.get('titulo', 'Resumen'))

            tf = sld.placeholders[1].text_frame

            pts = info.get('puntos', [])

            if pts:

                tf.text = str(pts[0])

                for pt in pts[1:]:

                    p = tf.add_paragraph()

                    p.text = str(pt)

        except: continue

    if df is not None:

        graficos = generar_graficos_matplotlib(df)

        for gbuf in graficos:

            sld = prs.slides.add_slide(prs.slide_layouts[6]) # blank slide

            sld.shapes.add_picture(gbuf, Inches(1), Inches(1), width=Inches(8))

    pptx_io = io.BytesIO()

    prs.save(pptx_io)

    return pptx_io.getvalue()

# ═══════════════════════════════════════════════════════════════

# [INYECCIÓN] Generador de Excel con openpyxl

# ═══════════════════════════════════════════════════════════════

def generar_excel(df: pd.DataFrame) -> bytes:

    """Exporta el DataFrame completo a un archivo Excel."""

    excel_io = io.BytesIO()

    with pd.ExcelWriter(excel_io, engine='openpyxl') as writer:

        df.to_excel(writer, index=False, sheet_name='Datos Analizados')

    return excel_io.getvalue()

@app.post("/upload")

async def upload_file(session_id: str = Form(...), file: UploadFile = File(...)):

    session = get_session(session_id)

    contents = await file.read()

    try:

        if file.filename.endswith(".csv"):

            session.df = pd.read_csv(io.BytesIO(contents))

        elif file.filename.endswith((".xls", ".xlsx")):

            session.df = pd.read_excel(io.BytesIO(contents))

        elif file.filename.endswith(".json"):

            session.df = pd.read_json(io.BytesIO(contents))

        elif file.filename.endswith(".txt"):

            text = contents.decode('utf-8', errors='replace')

            lines = [l.strip() for l in text.split('\n') if l.strip()]

            session.df = pd.DataFrame({'linea': range(1, len(lines)+1), 'contenido': lines})

        elif file.filename.endswith(".pdf"):

            from PyPDF2 import PdfReader

            reader = PdfReader(io.BytesIO(contents))

            pages_text = [page.extract_text() or '' for page in reader.pages]

            session.df = pd.DataFrame({'pagina': range(1, len(pages_text)+1), 'contenido': pages_text})

        elif file.filename.endswith(".docx"):

            doc_obj = Document(io.BytesIO(contents))

            paragraphs = [p.text for p in doc_obj.paragraphs if p.text.strip()]

            session.df = pd.DataFrame({'parrafo': range(1, len(paragraphs)+1), 'contenido': paragraphs})

        else:

            raise HTTPException(status_code=400, detail="Formato no soportado. Usa CSV, Excel, JSON, TXT, PDF o Word.")

        preview = session.df.head(5).to_dict(orient='records')

        return {

            "status": "success", 

            "filename": file.filename,

            "rows": len(session.df),

            "cols": len(session.df.columns),

            "columns": session.df.columns.tolist(),

            "preview": preview

        }

    except Exception as e:

        raise HTTPException(status_code=500, detail=str(e))

# ═══════════════════════════════════════════════════════════════

# [INYECCIÓN] Endpoint de chat con STREAMING (SSE)

# ═══════════════════════════════════════════════════════════════

@app.post("/chat/stream")

async def chat_stream(request: ChatRequest):

    session = get_session(request.session_id)

    if session.df is None:

        raise HTTPException(status_code=400, detail="No data uploaded for this session")

    llm = get_llm(request.model_choice)

    if llm is None:

        raise HTTPException(status_code=400, detail="Model API Key missing or invalid")

    session.messages.append({"role": "user", "content": request.prompt})

    try:

        agent = create_pandas_dataframe_agent(llm, session.df, verbose=False, allow_dangerous_code=True, agent_type="tool-calling")

        respuesta = agent.invoke(request.prompt)

        raw_output = respuesta.get('output', str(respuesta))

        output_text = str(raw_output)

        session.messages.append({"role": "assistant", "content": output_text})

        # Generar documentos en background con gráficos reales

        session.word_bytes = generar_informe_word(output_text, session.df)

        session.pdf_bytes = generar_informe_pdf(output_text, session.df)

        session.pptx_bytes = generar_pptx_presentacion(output_text, llm, session.df)

        session.excel_bytes = generar_excel(session.df)

    except Exception as e:

        output_text = f"Error en el análisis: {str(e)}"

        session.messages.append({"role": "assistant", "content": output_text})

    b64_images = []

    if session.df is not None:

        b64_images = [base64.b64encode(g.getvalue()).decode('utf-8') for g in generar_graficos_matplotlib(session.df)]

    # ═══ [INYECCIÓN] Streaming SSE: enviar letra por letra ═══

    async def generate():

        for char in output_text:

            yield {"data": json.dumps({"char": char})}

            await asyncio.sleep(0.008)  # 8ms por caracter = efecto typewriter

        yield {"data": json.dumps({"done": True, "images": b64_images})}

    return EventSourceResponse(generate())

# Endpoint clásico (fallback sin streaming)

@app.post("/chat")

async def chat_with_agent(request: ChatRequest):

    session = get_session(request.session_id)

    if session.df is None:

        raise HTTPException(status_code=400, detail="No data uploaded for this session")

    llm = get_llm(request.model_choice)

    if llm is None:

        raise HTTPException(status_code=400, detail="Model API Key missing or invalid")

    session.messages.append({"role": "user", "content": request.prompt})

    try:

        agent = create_pandas_dataframe_agent(llm, session.df, verbose=False, allow_dangerous_code=True, agent_type="tool-calling")

        respuesta = agent.invoke(request.prompt)

        raw_output = respuesta.get('output', str(respuesta))

        output_text = str(raw_output)

        session.messages.append({"role": "assistant", "content": output_text})

        # Generar documentos con gráficos reales

        session.word_bytes = generar_informe_word(output_text, session.df)

        session.pdf_bytes = generar_informe_pdf(output_text, session.df)

        session.pptx_bytes = generar_pptx_presentacion(output_text, llm, session.df)

        session.excel_bytes = generar_excel(session.df)

        b64_images = []

        if session.df is not None:

            b64_images = [base64.b64encode(g.getvalue()).decode('utf-8') for g in generar_graficos_matplotlib(session.df)]

        return {"response": output_text, "images": b64_images}

    except Exception as e:

        raise HTTPException(status_code=500, detail=str(e))

# ═══════════════════════════════════════════════════════════════

# [INYECCIÓN] Upload de audio/imagen multimedia

# ═══════════════════════════════════════════════════════════════

@app.post("/chat/multimedia")

async def chat_multimedia(

    session_id: str = Form(...),

    model_choice: str = Form(...),

    prompt: str = Form("Analiza este archivo adjunto."),

    media: UploadFile = File(...)

):

    session = get_session(session_id)

    if session.df is None:

        raise HTTPException(status_code=400, detail="No data uploaded")

    llm = get_llm(model_choice)

    if llm is None:

        raise HTTPException(status_code=400, detail="Model API Key missing")

    media_bytes = await media.read()

    output_text = ""

    try:

        # Generar data URI para devolver al frontend

        media_b64 = base64.b64encode(media_bytes).decode('utf-8')

        media_data_uri = f"data:{media.content_type};base64,{media_b64}"

        if media.content_type.startswith("audio/"):

            # 1. Transcribir Audio usando Groq Whisper

            groq_key_live = os.getenv("GROQ_API_KEY")

            if groq_key_live:

                ext = "webm" if "webm" in media.content_type else "mp3"

                fname = f"audio.{ext}"

                try:

                    response = httpx.post(

                        "https://api.groq.com/openai/v1/audio/transcriptions",

                        headers={"Authorization": f"Bearer {groq_key_live}"},

                        files={"file": (fname, media_bytes, media.content_type)},

                        data={"model": "whisper-large-v3", "response_format": "json"}

                    )

                    transcription = response.json().get("text", "")

                    if transcription:

                        full_prompt = f"El usuario envió una nota de voz diciendo: '{transcription}'. Responde a su solicitud basándote en los datos."

                    else:

                        full_prompt = "El usuario envió un audio pero no se pudo transcribir. Pregúntale si puede escribirlo."

                except Exception as ex:

                    full_prompt = f"Error al transcribir audio: {str(ex)}"

            else:

                full_prompt = "No hay API Key de Groq para transcribir el audio."

            session.messages.append({"role": "user", "content": f"🎤 Audio transcrito: {full_prompt}"})

            agent = create_pandas_dataframe_agent(llm, session.df, verbose=False, allow_dangerous_code=True, agent_type="tool-calling")

            respuesta = agent.invoke(full_prompt)

            output_text = str(respuesta.get('output', str(respuesta)))

        elif media.content_type.startswith("image/"):

            # 2. Analizar Imagen directamente con Gemini Multimodal

            google_key_live = os.getenv("GOOGLE_API_KEY")

            if google_key_live:

                b64_img = base64.b64encode(media_bytes).decode('utf-8')

                gemini_vision = ChatGoogleGenerativeAI(model="gemini-1.5-flash", temperature=0.1, google_api_key=google_key_live)

                from langchain_core.messages import HumanMessage

                msg = HumanMessage(

                    content=[

                        {"type": "text", "text": "Analiza esta imagen relacionada con mi análisis de datos. ¿Qué ves y qué recomendaciones me das?"},

                        {"type": "image_url", "image_url": {"url": f"data:{media.content_type};base64,{b64_img}"}}

                    ]

                )

                session.messages.append({"role": "user", "content": f"🖼️ Imagen adjunta: {media.filename}"})

                res_vision = gemini_vision.invoke([msg])

                output_text = "Análisis de la imagen:\n" + str(res_vision.content)

            else:

                raise HTTPException(status_code=400, detail="Se requiere GOOGLE_API_KEY para analizar imágenes.")

        else:

            # Fallback para otros archivos

            media_desc = f"[Archivo adjunto: {media.filename}, {len(media_bytes)} bytes, tipo: {media.content_type}]"

            full_prompt = f"{prompt}\n\n{media_desc}"

            session.messages.append({"role": "user", "content": full_prompt})

            agent = create_pandas_dataframe_agent(llm, session.df, verbose=False, allow_dangerous_code=True, agent_type="tool-calling")

            respuesta = agent.invoke(full_prompt)

            output_text = str(respuesta.get('output', str(respuesta)))

        session.messages.append({"role": "assistant", "content": output_text})

        # Generar documentos actualizados si es posible

        if session.df is not None:

            session.word_bytes = generar_informe_word(output_text, session.df)

            session.pdf_bytes = generar_informe_pdf(output_text, session.df)

            session.pptx_bytes = generar_pptx_presentacion(output_text, llm, session.df)

            session.excel_bytes = generar_excel(session.df)

        b64_images = []

        if session.df is not None:

            b64_images = [base64.b64encode(g.getvalue()).decode('utf-8') for g in generar_graficos_matplotlib(session.df)]

        # Devolver media_type y media_url para renderizar en el chat

        result = {"response": output_text, "images": b64_images}

        if media.content_type.startswith("image/"):

            result["media_type"] = "image"

            result["media_url"] = media_data_uri

        # Audio ya lo tiene el frontend localmente, no reenviarlo por base64

        return result

    except Exception as e:

        import traceback

        traceback.print_exc()

        raise HTTPException(status_code=500, detail=str(e))

@app.get("/export/word/{session_id}/{filename}")

async def export_word(session_id: str, filename: str):

    session = get_session(session_id)

    if not session.word_bytes:

        raise HTTPException(status_code=404, detail="No report generated")

    return Response(content=session.word_bytes, media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document", headers={"Content-Disposition": f'attachment; filename="{filename}"'})

@app.get("/export/pdf/{session_id}/{filename}")

async def export_pdf(session_id: str, filename: str):

    session = get_session(session_id)

    if not session.pdf_bytes:

        raise HTTPException(status_code=404, detail="No report generated")

    return Response(content=session.pdf_bytes, media_type="application/pdf", headers={"Content-Disposition": f'attachment; filename="{filename}"'})

@app.get("/export/pptx/{session_id}/{filename}")

async def export_pptx(session_id: str, filename: str):

    session = get_session(session_id)

    if not session.pptx_bytes:

        raise HTTPException(status_code=404, detail="No report generated")

    return Response(content=session.pptx_bytes, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", headers={"Content-Disposition": f'attachment; filename="{filename}"'})

# ═══════════════════════════════════════════════════════════════

# [INYECCIÓN] Endpoint de exportación Excel

# ═══════════════════════════════════════════════════════════════

@app.get("/export/excel/{session_id}/{filename}")

async def export_excel(session_id: str, filename: str):

    session = get_session(session_id)

    if session.df is None:

        raise HTTPException(status_code=404, detail="No data loaded")

    # Generar Excel al momento si no existe

    if not session.excel_bytes:

        session.excel_bytes = generar_excel(session.df)

    return Response(

        content=session.excel_bytes,

        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",

        headers={"Content-Disposition": f'attachment; filename="{filename}"'}

    )

@app.post("/chat/help")

async def chat_help(request: ChatRequest):

    llm = get_llm(request.model_choice)

    if llm is None:

        return {"response": "Configura una API Key válida primero para usar el bot de manera inteligente. Por ahora respondo tus dudas en base a mis opciones límite."}

    prompt = f"Eres NexusBot, un asistente ejecutivo muy servicial que navega por Nexus Analytics. Responde de forma clara y directa a esta duda del usuario relacionada con análisis de datos o la plataforma: {request.prompt}."

    try:

        from langchain_core.messages import HumanMessage

        res = llm.invoke([HumanMessage(content=prompt)])

        return {"response": res.content}

    except Exception as e:

        return {"response": f"Lo siento, estoy teniendo problemas de red: {str(e)}"}

